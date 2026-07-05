import json
from collections import Counter, defaultdict
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "PSI_Transition_State_Project_Detailed_Presentation.pptx"
ASSET_DIR = ROOT / "ppt_assets"

TITLE = "PSI Transition-State Prediction"
SUBTITLE = "Deep learning pipeline for predicting 3D transition-state geometry and activation energy from reactant/product structures"

COLORS = {
    "ink": RGBColor(24, 31, 42),
    "muted": RGBColor(84, 96, 112),
    "line": RGBColor(210, 217, 226),
    "bg": RGBColor(247, 249, 252),
    "blue": RGBColor(42, 99, 181),
    "teal": RGBColor(0, 137, 123),
    "green": RGBColor(46, 125, 50),
    "orange": RGBColor(217, 119, 6),
    "red": RGBColor(185, 28, 28),
    "purple": RGBColor(109, 40, 217),
    "white": RGBColor(255, 255, 255),
}


def load_json(name):
    with open(ROOT / name, "r", encoding="utf-8") as f:
        return json.load(f)


def safe_load_checkpoint_metadata():
    try:
        import torch

        ckpt = torch.load(ROOT / "psi_final.pt", map_location="cpu", weights_only=False)
        return ckpt.get("metadata", {})
    except Exception:
        return {}


def regression_metrics(rows, key="Ea_pred"):
    y = np.array([r["Ea_true"] for r in rows], dtype=float)
    p = np.array([r[key] for r in rows], dtype=float)
    err = p - y
    abs_err = np.abs(err)
    ss_tot = float(np.sum((y - y.mean()) ** 2))
    r2 = 1.0 - float(np.sum(err**2)) / ss_tot if ss_tot else 0.0
    corr = float(np.corrcoef(y, p)[0, 1]) if len(y) > 1 else 0.0
    return {
        "count": len(rows),
        "mae": float(abs_err.mean()) if len(rows) else 0.0,
        "rmse": float(np.sqrt(np.mean(err**2))) if len(rows) else 0.0,
        "r2": r2,
        "pearson": corr,
        "median_abs_error": float(np.median(abs_err)) if len(rows) else 0.0,
        "p90_abs_error": float(np.percentile(abs_err, 90)) if len(rows) else 0.0,
    }


def geometry_metrics(rows):
    d = np.array([r["dist_MAE"] for r in rows], dtype=float)
    atoms = np.array([r["n_atoms"] for r in rows], dtype=float)
    return {
        "dist_mae": float(d.mean()) if len(d) else 0.0,
        "dist_std": float(d.std()) if len(d) else 0.0,
        "dist_median": float(np.median(d)) if len(d) else 0.0,
        "dist_p90": float(np.percentile(d, 90)) if len(d) else 0.0,
        "atoms_mean": float(atoms.mean()) if len(atoms) else 0.0,
        "atoms_min": int(atoms.min()) if len(atoms) else 0,
        "atoms_max": int(atoms.max()) if len(atoms) else 0,
    }


def summarize_raw_dataset():
    try:
        raw = load_json("extracted_dataset.json")
    except FileNotFoundError:
        return {}

    roles = Counter()
    rxn_roles = defaultdict(set)
    atom_counts = Counter()
    for entry in raw:
        parts = entry.get("filename", "").split("/")
        if len(parts) >= 3:
            rxn_id = parts[1]
            prefix = parts[2].lower()
            role = (
                "reactant"
                if prefix.startswith("r")
                else "product"
                if prefix.startswith("p")
                else "transition state"
                if prefix.startswith("ts")
                else "unknown"
            )
            roles[role] += 1
            rxn_roles[rxn_id].add(role)
        for atom in entry.get("atoms", []):
            atom_counts[atom.get("atom", "?")] += 1

    complete = sum(
        1 for role_set in rxn_roles.values()
        if {"reactant", "product", "transition state"}.issubset(role_set)
    )
    return {
        "raw_entries": len(raw),
        "unique_reactions": len(rxn_roles),
        "complete_triplets": complete,
        "roles": dict(roles),
        "atom_counts": dict(atom_counts),
    }


def collect_metrics():
    data = load_json("detailed_analysis.json")
    history = load_json("training_history.json")
    metadata = safe_load_checkpoint_metadata()
    config = metadata.get("config_snapshot", {})
    raw_summary = summarize_raw_dataset()

    train = [r for r in data if r["split"] == "train"]
    val = [r for r in data if r["split"] == "val"]
    all_rows = data

    metrics = {
        "train": {**regression_metrics(train), **geometry_metrics(train)},
        "val": {**regression_metrics(val), **geometry_metrics(val)},
        "all": {**regression_metrics(all_rows), **geometry_metrics(all_rows)},
        "physics_train": regression_metrics(train, "Ea_pred_physics"),
        "physics_val": regression_metrics(val, "Ea_pred_physics"),
        "physics_all": regression_metrics(all_rows, "Ea_pred_physics"),
        "atom_types": sorted({a for r in data for a in r["atom_types"]}),
        "atom_count_distribution": Counter(r["n_atoms"] for r in data),
        "energy_range": {
            "true_min": min(r["Ea_true"] for r in data),
            "true_max": max(r["Ea_true"] for r in data),
            "pred_min": min(r["Ea_pred"] for r in data),
            "pred_max": max(r["Ea_pred"] for r in data),
        },
        "history": history,
        "config": config,
        "metadata": metadata,
        "raw_summary": raw_summary,
        "data": data,
    }
    if history:
        metrics["best_val_select"] = min(history, key=lambda h: h.get("val_select", h["val_loss"]))
        metrics["best_val_geom"] = min(history, key=lambda h: h["val_geom"])
        metrics["best_val_ea"] = min(history, key=lambda h: h["val_ea_mae"])
    return metrics


def ensure_assets():
    ASSET_DIR.mkdir(exist_ok=True)


def set_font(run, size=18, bold=False, color=None):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_bg(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, kicker=None):
    if kicker:
        box = slide.shapes.add_textbox(Inches(0.65), Inches(0.28), Inches(11.0), Inches(0.28))
        p = box.text_frame.paragraphs[0]
        p.text = kicker.upper()
        p.font.name = "Aptos"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLORS["teal"]

    box = slide.shapes.add_textbox(Inches(0.62), Inches(0.52), Inches(12.0), Inches(0.62))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(29)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    add_rule(slide, y=1.18)


def add_rule(slide, y=1.2):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(y), Inches(12.0), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.15), Inches(11.8), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"PSI Transition-State Prediction | {idx}"
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["muted"]
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items, x, y, w, h, font_size=16, color=COLORS["ink"], gap=4):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return box


def add_number_card(slide, x, y, w, h, value, label, color=COLORS["blue"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    set_font(r, 23, True, color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    p2.font.name = "Aptos"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLORS["muted"]
    return shape


def add_section_label(slide, text, x, y, w, color=COLORS["teal"]):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = color
    return box


def add_table(slide, x, y, w, h, headers, rows, font_size=10.5):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["ink"]
        for p in cell.text_frame.paragraphs:
            p.font.name = "Aptos"
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = COLORS["white"]
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 else RGBColor(242, 246, 250)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["ink"]
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
    return table_shape


def add_flow_box(slide, x, y, w, h, title, body, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["white"]
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 12, True, color)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLORS["muted"]
    return box


def add_arrow(slide, x, y, w=0.36):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.22))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["line"]
    arrow.line.fill.background()
    return arrow


def add_callout(slide, x, y, w, h, text, color=COLORS["blue"], font_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(238, 244, 255)
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    return shape


def add_code_box(slide, x, y, w, h, text, font_size=10.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(28, 36, 49)
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Consolas"
    p.font.size = Pt(font_size)
    p.font.color.rgb = COLORS["white"]
    return box


def metric_cfg(metrics, key, default):
    value = metrics.get("config", {}).get(key, default)
    return default if value is None else value


def chart_style():
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "axes.edgecolor": "#cbd5e1",
        "axes.labelcolor": "#1f2937",
        "xtick.color": "#475569",
        "ytick.color": "#475569",
        "figure.facecolor": "white",
        "axes.facecolor": "white",
    })


def save_training_chart(metrics):
    chart_style()
    hist = metrics["history"]
    epochs = [h["epoch"] for h in hist]
    fig, ax1 = plt.subplots(figsize=(9.5, 4.9), dpi=170)
    ax1.plot(epochs, [h["train_geom"] for h in hist], label="Train geometry loss", color="#00897b", linewidth=1.8)
    ax1.plot(epochs, [h["val_geom"] for h in hist], label="Val geometry loss", color="#2a63b5", linewidth=1.8)
    ax1.set_xlabel("Epoch")
    ax1.set_ylabel("Geometry loss")
    ax1.grid(True, alpha=0.25)
    ax2 = ax1.twinx()
    ax2.plot(epochs, [h["val_ea_mae"] for h in hist], label="Val Ea MAE", color="#d97706", linewidth=1.6)
    ax2.set_ylabel("Ea MAE (kcal/mol)")
    ax1.axvline(metrics["config"].get("ea_warmup_epochs", 200), color="#64748b", linestyle="--", linewidth=1)
    ax1.text(metrics["config"].get("ea_warmup_epochs", 200) + 12, ax1.get_ylim()[1] * 0.86, "Ea head warmup ends", color="#64748b", fontsize=8)
    lines = ax1.get_lines() + ax2.get_lines()
    ax1.legend(lines, [line.get_label() for line in lines], loc="upper right", frameon=False)
    fig.tight_layout()
    path = ASSET_DIR / "training_curves.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_parity_chart(metrics):
    chart_style()
    data = metrics["data"]
    fig, ax = plt.subplots(figsize=(6.5, 5.6), dpi=180)
    for split, color, label, alpha in [("train", "#00897b", "Train", 0.28), ("val", "#d97706", "Validation", 0.72)]:
        rows = [r for r in data if r["split"] == split]
        ax.scatter(
            [r["Ea_true"] for r in rows],
            [r["Ea_pred"] for r in rows],
            s=9 if split == "train" else 14,
            alpha=alpha,
            color=color,
            edgecolors="none",
            label=label,
        )
    lo = min(metrics["energy_range"]["true_min"], metrics["energy_range"]["pred_min"]) - 3
    hi = max(metrics["energy_range"]["true_max"], metrics["energy_range"]["pred_max"]) + 3
    ax.plot([lo, hi], [lo, hi], color="#1f2937", linewidth=1.1, linestyle="--", label="Ideal")
    ax.set_xlim(lo, hi)
    ax.set_ylim(lo, hi)
    ax.set_xlabel("True Ea (kcal/mol)")
    ax.set_ylabel("Predicted Ea (kcal/mol)")
    ax.set_title("Activation Energy Parity")
    ax.grid(True, alpha=0.25)
    ax.legend(frameon=False, loc="upper left")
    fig.tight_layout()
    path = ASSET_DIR / "ea_parity.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_error_hist_chart(metrics):
    chart_style()
    data = metrics["data"]
    fig, ax = plt.subplots(figsize=(7.3, 4.7), dpi=180)
    for split, color in [("train", "#00897b"), ("val", "#d97706")]:
        rows = [r for r in data if r["split"] == split]
        errors = [abs(r["Ea_pred"] - r["Ea_true"]) for r in rows]
        ax.hist(errors, bins=45, alpha=0.55, color=color, label=split.capitalize())
    ax.set_xlabel("Absolute Ea error (kcal/mol)")
    ax.set_ylabel("Reaction count")
    ax.set_title("Error Distribution")
    ax.grid(True, axis="y", alpha=0.25)
    ax.legend(frameon=False)
    fig.tight_layout()
    path = ASSET_DIR / "ea_error_hist.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_distance_chart(metrics):
    chart_style()
    data = metrics["data"]
    fig, ax = plt.subplots(figsize=(7.3, 4.7), dpi=180)
    for split, color in [("train", "#2a63b5"), ("val", "#b91c1c")]:
        rows = [r for r in data if r["split"] == split]
        ax.hist([r["dist_MAE"] for r in rows], bins=45, alpha=0.55, color=color, label=split.capitalize())
    ax.set_xlabel("Distance-matrix MAE (Angstrom)")
    ax.set_ylabel("Reaction count")
    ax.set_title("Geometry Error Distribution")
    ax.grid(True, axis="y", alpha=0.25)
    ax.legend(frameon=False)
    fig.tight_layout()
    path = ASSET_DIR / "dist_error_hist.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_baseline_chart(metrics):
    chart_style()
    labels = ["Train", "Validation", "All"]
    neural = [metrics["train"]["mae"], metrics["val"]["mae"], metrics["all"]["mae"]]
    physics = [metrics["physics_train"]["mae"], metrics["physics_val"]["mae"], metrics["physics_all"]["mae"]]
    x = np.arange(len(labels))
    width = 0.34
    fig, ax = plt.subplots(figsize=(7.3, 4.4), dpi=180)
    ax.bar(x - width / 2, neural, width, label="Neural Ea head", color="#00897b")
    ax.bar(x + width / 2, physics, width, label="Physics baseline", color="#64748b")
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylabel("Ea MAE (kcal/mol)")
    ax.set_title("Activation Energy Model Comparison")
    ax.grid(True, axis="y", alpha=0.25)
    ax.legend(frameon=False)
    for i, val in enumerate(neural):
        ax.text(i - width / 2, val + 0.35, f"{val:.1f}", ha="center", fontsize=8)
    for i, val in enumerate(physics):
        ax.text(i + width / 2, val + 0.35, f"{val:.1f}", ha="center", fontsize=8)
    fig.tight_layout()
    path = ASSET_DIR / "baseline_comparison.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_atom_distribution_chart(metrics):
    chart_style()
    dist = metrics["atom_count_distribution"]
    xs = sorted(dist)
    ys = [dist[x] for x in xs]
    fig, ax = plt.subplots(figsize=(7.2, 4.3), dpi=180)
    ax.bar(xs, ys, color="#2a63b5")
    ax.set_xlabel("Atoms per reaction")
    ax.set_ylabel("Reaction count")
    ax.set_title("Molecular Size Distribution")
    ax.grid(True, axis="y", alpha=0.25)
    ax.set_xticks(xs)
    fig.tight_layout()
    path = ASSET_DIR / "atom_distribution.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def generate_charts(metrics):
    ensure_assets()
    return {
        "training": save_training_chart(metrics),
        "parity": save_parity_chart(metrics),
        "errors": save_error_hist_chart(metrics),
        "distance": save_distance_chart(metrics),
        "baseline": save_baseline_chart(metrics),
        "atoms": save_atom_distribution_chart(metrics),
    }


def init_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    return slide


def fmt_int(value):
    return f"{int(value):,}"


def fmt_float(value, digits=2):
    return f"{float(value):.{digits}f}"


def add_deck_slides(prs, metrics, charts):
    slide_no = 1

    # Slide 1
    slide = blank_slide(prs)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(236, 244, 250)
    shape.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.23))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["teal"]
    band.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = "Aptos Display"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    sub = slide.shapes.add_textbox(Inches(0.85), Inches(2.43), Inches(10.8), Inches(0.85))
    p = sub.text_frame.paragraphs[0]
    p.text = SUBTITLE
    p.font.name = "Aptos"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["muted"]
    add_number_card(slide, 0.85, 4.1, 2.15, 0.95, fmt_int(metrics["all"]["count"]), "evaluated reactions", COLORS["blue"])
    add_number_card(slide, 3.25, 4.1, 2.15, 0.95, f"{metrics['val']['mae']:.2f}", "val Ea MAE kcal/mol", COLORS["teal"])
    add_number_card(slide, 5.65, 4.1, 2.15, 0.95, f"{metrics['val']['dist_mae']:.4f}", "val distance MAE A", COLORS["orange"])
    add_number_card(slide, 8.05, 4.1, 2.15, 0.95, f"{metrics['val']['r2']:.3f}", "val Ea R2", COLORS["purple"])
    add_callout(
        slide,
        0.85,
        5.55,
        11.1,
        0.7,
        "Input: reactant + product Gaussian geometries. Output: predicted transition-state distance matrix, 3D coordinates, and activation energy.",
        COLORS["blue"],
        14,
    )
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 2
    slide = blank_slide(prs)
    add_title(slide, "Project Motivation", "Context")
    add_bullets(
        slide,
        [
            "Transition states control reaction rates, but TS optimization is expensive and sensitive to initial guesses.",
            "The project learns the mapping from endpoint structures to a transition-state geometry using only reactant and product geometries.",
            "A second objective predicts activation energy (Ea), enabling both structural and kinetic screening.",
            "The implementation targets a practical pipeline: data extraction, model training, prediction, and an interactive dashboard.",
        ],
        0.75,
        1.55,
        6.25,
        4.5,
        17,
    )
    add_flow_box(slide, 7.35, 1.7, 4.6, 0.78, "Scientific question", "Can endpoint molecular geometries carry enough signal to estimate a TS?", COLORS["teal"])
    add_flow_box(slide, 7.35, 2.75, 4.6, 0.78, "Engineering question", "Can that estimate be packaged into a repeatable train/predict/dashboard workflow?", COLORS["blue"])
    add_flow_box(slide, 7.35, 3.8, 4.6, 0.78, "Output value", "Fast candidate TS geometry and barrier estimates for downstream analysis.", COLORS["orange"])
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 3
    slide = blank_slide(prs)
    add_title(slide, "Project Scope and Deliverables", "Overview")
    rows = [
        ["Data extraction", "Parse Gaussian logs into JSON with atom types, coordinates, and energies"],
        ["Training pipeline", "Build reaction triplets, normalize features, train neural geometry + Ea model"],
        ["Prediction CLI", "Predict TS geometry and Ea from new reactant/product log files"],
        ["Visualization", "Generate an HTML dashboard with parity plots, error histograms, and 3D cases"],
        ["Model artifacts", "Save best/final PyTorch checkpoints with normalization metadata"],
    ]
    add_table(slide, 0.8, 1.55, 11.65, 3.2, ["Deliverable", "What it does"], rows, 12)
    add_callout(slide, 1.2, 5.35, 10.6, 0.68, "Main artifacts in this workspace: psi_cloud_pipeline.py, psi_final.pt, detailed_analysis.json, training_history.json, psi_results_dashboard.html.", COLORS["teal"], 13)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 4
    slide = blank_slide(prs)
    add_title(slide, "Dataset Summary", "Data")
    raw = metrics["raw_summary"]
    add_number_card(slide, 0.8, 1.55, 2.0, 0.86, fmt_int(raw.get("raw_entries", 0)), "raw JSON entries", COLORS["blue"])
    add_number_card(slide, 3.05, 1.55, 2.0, 0.86, fmt_int(raw.get("unique_reactions", 0)), "unique reactions", COLORS["teal"])
    add_number_card(slide, 5.3, 1.55, 2.0, 0.86, fmt_int(raw.get("complete_triplets", 0)), "complete triplets", COLORS["green"])
    add_number_card(slide, 7.55, 1.55, 2.0, 0.86, fmt_int(metrics["all"]["count"]), "used reactions", COLORS["orange"])
    add_number_card(slide, 9.8, 1.55, 2.0, 0.86, "/".join(metrics["atom_types"]), "atom types", COLORS["purple"])
    slide.shapes.add_picture(str(charts["atoms"]), Inches(0.85), Inches(2.85), Inches(5.95), Inches(3.55))
    rows = [
        ["Train split", fmt_int(metrics["train"]["count"]), f"{metrics['train']['atoms_mean']:.1f}", f"{metrics['train']['atoms_min']}-{metrics['train']['atoms_max']}"],
        ["Validation split", fmt_int(metrics["val"]["count"]), f"{metrics['val']['atoms_mean']:.1f}", f"{metrics['val']['atoms_min']}-{metrics['val']['atoms_max']}"],
        ["All", fmt_int(metrics["all"]["count"]), f"{metrics['all']['atoms_mean']:.1f}", f"{metrics['all']['atoms_min']}-{metrics['all']['atoms_max']}"],
    ]
    add_table(slide, 7.25, 3.0, 4.95, 1.35, ["Split", "Count", "Avg atoms", "Range"], rows, 10.5)
    atom_counts = raw.get("atom_counts", {})
    atom_text = "Raw atom occurrences: " + ", ".join(f"{k}={fmt_int(v)}" for k, v in sorted(atom_counts.items()))
    add_callout(slide, 7.25, 4.8, 4.95, 0.72, atom_text, COLORS["blue"], 11)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 5
    slide = blank_slide(prs)
    add_title(slide, "End-to-End Pipeline", "Workflow")
    x0, y0, w, h, gap = 0.62, 2.05, 1.75, 1.08, 0.24
    boxes = [
        ("Gaussian logs", "R / P / TS structures and energies", COLORS["blue"]),
        ("Extraction", "Parse atoms, coordinates, SCF energies", COLORS["teal"]),
        ("Triplets", "Build complete reaction samples", COLORS["green"]),
        ("Features", "Distances, masks, atom descriptors", COLORS["orange"]),
        ("PSI model", "GRU + Transformer + geometry head", COLORS["purple"]),
        ("EGNN + Ea", "3D refinement and barrier head", COLORS["red"]),
    ]
    for i, (title, body, color) in enumerate(boxes):
        add_flow_box(slide, x0 + i * (w + gap), y0, w, h, title, body, color)
        if i < len(boxes) - 1:
            add_arrow(slide, x0 + (i + 1) * w + i * gap + 0.02, y0 + 0.43, 0.25)
    add_flow_box(slide, 2.25, 4.25, 2.05, 0.88, "Artifacts", "psi_best.pt, psi_final.pt, training_history.json", COLORS["blue"])
    add_flow_box(slide, 5.0, 4.25, 2.05, 0.88, "Evaluation", "detailed_analysis.json with geometry + Ea metrics", COLORS["teal"])
    add_flow_box(slide, 7.75, 4.25, 2.05, 0.88, "Dashboard", "Plotly/3Dmol HTML performance report", COLORS["orange"])
    add_arrow(slide, 4.48, 4.57, 0.34)
    add_arrow(slide, 7.23, 4.57, 0.34)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 6
    slide = blank_slide(prs)
    add_title(slide, "Data Processing and Feature Engineering", "Method")
    add_bullets(
        slide,
        [
            "Reaction IDs are grouped into reactant/product/TS triplets from filenames.",
            "Molecules above max_atoms are skipped; all coordinate and distance tensors are padded to a fixed 30-atom size.",
            "Reactant fragments are Kabsch-aligned to products before building the midpoint/interpolated geometry input.",
            "Per-atom physical descriptors are attached to nodes: electronegativity, atomic number, and atomic mass.",
            "Global energy features include signed reaction energy, composition counts, displacement statistics, and bond-angle statistics.",
            "Geometry masks focus the loss on atoms in connected TS fragments and avoid padding noise.",
        ],
        0.85,
        1.5,
        7.0,
        4.85,
        15,
    )
    rows = [
        ["D_R", "Reactant distance matrix"],
        ["D_I", "Aligned midpoint/interpolated distance matrix"],
        ["D_P", "Product distance matrix"],
        ["atom_ids", "Learned atom identity embedding"],
        ["atom_phys", "EN, Z, mass per atom"],
        ["energy_feats", "20D reaction/global descriptor"],
    ]
    add_table(slide, 8.15, 1.58, 4.1, 3.2, ["Tensor", "Meaning"], rows, 10.2)
    add_callout(slide, 8.15, 5.15, 4.1, 0.74, "Activation energy target in this run is TS energy minus the higher-energy endpoint, converted from Hartree to kcal/mol.", COLORS["orange"], 11.5)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 7
    slide = blank_slide(prs)
    add_title(slide, "Model Architecture", "Architecture")
    add_flow_box(slide, 0.85, 1.75, 2.05, 0.92, "Distance encoding", "32 Gaussian basis features over pairwise distances", COLORS["blue"])
    add_arrow(slide, 3.1, 2.1, 0.42)
    add_flow_box(slide, 3.62, 1.75, 2.05, 0.92, "PSICore", "2-layer bidirectional GRU across R/I/P states", COLORS["teal"])
    add_arrow(slide, 5.9, 2.1, 0.42)
    add_flow_box(slide, 6.42, 1.75, 2.05, 0.92, "Transformer", "3 pre-norm encoder layers, 8 attention heads", COLORS["green"])
    add_arrow(slide, 8.7, 2.1, 0.42)
    add_flow_box(slide, 9.22, 1.75, 2.15, 0.92, "Geometry head", "Pairwise alpha interpolation + learned delta", COLORS["orange"])

    add_flow_box(slide, 2.05, 3.75, 2.15, 0.9, "MDS", "Differentiable distance-to-coordinate embedding", COLORS["blue"])
    add_arrow(slide, 4.42, 4.07, 0.42)
    add_flow_box(slide, 4.95, 3.75, 2.15, 0.9, "EGNN", "4 equivariant graph convolution layers refine TS coords", COLORS["purple"])
    add_arrow(slide, 7.32, 4.07, 0.42)
    add_flow_box(slide, 7.85, 3.75, 2.15, 0.9, "Ea head", "Masked mean pooled node features + reaction energy", COLORS["red"])
    add_callout(slide, 1.1, 5.7, 10.95, 0.66, "The model predicts a TS distance matrix first, refines it as 3D coordinates, then uses the refined local environment to predict Ea.", COLORS["teal"], 13)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 8
    slide = blank_slide(prs)
    add_title(slide, "Geometry Prediction Objective", "Modeling")
    add_bullets(
        slide,
        [
            "The geometry head predicts each TS pairwise distance from learned atom context, endpoint distances, and atom descriptors.",
            "It starts from a constrained base: alpha * D_R + (1 - alpha) * D_P, then adds a clamped learned delta.",
            "Loss combines final geometry error, coarse pre-EGNN geometry error, and a spectator constraint for near-unchanged atom pairs.",
            "Post-processing symmetrizes the distance matrix, zeros the diagonal, clamps steric collisions, and recovers coordinates with fragment-aware MDS.",
        ],
        0.85,
        1.55,
        6.15,
        4.75,
        16,
    )
    add_section_label(slide, "Core formula", 7.35, 1.6, 4.2)
    formula = slide.shapes.add_textbox(Inches(7.35), Inches(2.0), Inches(4.75), Inches(0.95))
    p = formula.text_frame.paragraphs[0]
    p.text = "D_TS = alpha * D_R + (1 - alpha) * D_P + delta"
    p.font.name = "Aptos"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    add_flow_box(slide, 7.35, 3.25, 4.1, 0.86, "Why distance matrices?", "They are invariant to global translation and rotation, simplifying the geometry learning target.", COLORS["blue"])
    add_flow_box(slide, 7.35, 4.35, 4.1, 0.86, "Why EGNN refinement?", "The coordinate stage can adjust 3D placement while preserving equivariance to rigid transformations.", COLORS["purple"])
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 9
    slide = blank_slide(prs)
    add_title(slide, "Activation Energy Strategy", "Modeling")
    add_bullets(
        slide,
        [
            "The learned Ea head trains jointly with geometry after a warmup period so it receives useful predicted TS structures.",
            "Input to the head: EGNN-refined per-atom node features, signed reaction energy, and the 20D energy descriptor vector.",
            "Output is normalized Ea; inference denormalizes using checkpoint metadata.",
            "A physics baseline is retained for comparison using Marcus-style reorganization energy, Hammond index, signed reaction energy, and OLS.",
        ],
        0.85,
        1.55,
        6.75,
        4.5,
        16,
    )
    rows = [
        ["Learned Ea head", "Primary reported prediction", f"Val MAE {metrics['val']['mae']:.2f}"],
        ["Physics baseline", "Interpretability/comparison path", f"Val MAE {metrics['physics_val']['mae']:.2f}"],
        ["Warmup", "Avoids early Ea gradients from poor geometries", f"{metrics['config'].get('ea_warmup_epochs', 200)} epochs"],
    ]
    add_table(slide, 8.0, 1.8, 4.35, 1.85, ["Component", "Purpose", "Result"], rows, 10.2)
    add_callout(slide, 8.0, 4.35, 4.35, 0.85, "In this saved checkpoint metadata, the neural Ea head is the active prediction source in detailed_analysis.json.", COLORS["teal"], 11.5)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 10
    slide = blank_slide(prs)
    add_title(slide, "Training Configuration", "Training")
    cfg = metrics["config"]
    rows = [
        ["Target reactions", fmt_int(cfg.get("target_reactions", metrics["all"]["count"]))],
        ["Train / validation split", f"{metrics['train']['count']:,} / {metrics['val']['count']:,}"],
        ["Max atoms", cfg.get("max_atoms", 30)],
        ["Batch size", cfg.get("batch_size", 32)],
        ["Learning rate", cfg.get("lr", 1.5e-4)],
        ["Weight decay", cfg.get("weight_decay", 1e-2)],
        ["Warmup epochs", cfg.get("warmup_epochs", 40)],
        ["Patience", cfg.get("patience", 120)],
        ["AMP", cfg.get("amp", True)],
    ]
    add_table(slide, 0.9, 1.48, 4.5, 4.45, ["Setting", "Value"], rows, 10.3)
    add_bullets(
        slide,
        [
            "Optimizer: AdamW with cosine annealing and warmup.",
            "Checkpoint selection uses validation geometry, then includes normalized Ea after the Ea warmup.",
            "Training history contains 1,053 epochs; run stopped before the configured 1,500-epoch maximum.",
            "Best validation Ea epoch: "
            f"{metrics.get('best_val_ea', {}).get('epoch', 'n/a')} "
            f"({metrics.get('best_val_ea', {}).get('val_ea_mae', 0):.2f} kcal/mol).",
        ],
        6.05,
        1.65,
        5.8,
        3.2,
        15,
    )
    add_callout(slide, 6.05, 5.2, 5.4, 0.72, f"Final validation geometry loss: {metrics['history'][-1]['val_geom']:.5f}; final validation Ea MAE: {metrics['history'][-1]['val_ea_mae']:.2f} kcal/mol.", COLORS["blue"], 12)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 11
    slide = blank_slide(prs)
    add_title(slide, "Training Dynamics", "Results")
    slide.shapes.add_picture(str(charts["training"]), Inches(0.8), Inches(1.45), Inches(7.3), Inches(4.3))
    add_bullets(
        slide,
        [
            "Geometry loss improves rapidly early and continues improving slowly through long training.",
            f"Ea head activates after epoch {metrics['config'].get('ea_warmup_epochs', 200)}, after the geometry backbone has stabilized.",
            f"Best validation geometry in history: epoch {metrics['best_val_geom']['epoch']} with loss {metrics['best_val_geom']['val_geom']:.5f}.",
            f"Best validation Ea in history: epoch {metrics['best_val_ea']['epoch']} with MAE {metrics['best_val_ea']['val_ea_mae']:.2f} kcal/mol.",
        ],
        8.35,
        1.7,
        4.0,
        4.2,
        13.5,
    )
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 12
    slide = blank_slide(prs)
    add_title(slide, "Evaluation Metrics", "Results")
    rows = [
        ["Train", fmt_int(metrics["train"]["count"]), f"{metrics['train']['mae']:.2f}", f"{metrics['train']['rmse']:.2f}", f"{metrics['train']['r2']:.3f}", f"{metrics['train']['pearson']:.3f}", f"{metrics['train']['dist_mae']:.4f}"],
        ["Validation", fmt_int(metrics["val"]["count"]), f"{metrics['val']['mae']:.2f}", f"{metrics['val']['rmse']:.2f}", f"{metrics['val']['r2']:.3f}", f"{metrics['val']['pearson']:.3f}", f"{metrics['val']['dist_mae']:.4f}"],
        ["All", fmt_int(metrics["all"]["count"]), f"{metrics['all']['mae']:.2f}", f"{metrics['all']['rmse']:.2f}", f"{metrics['all']['r2']:.3f}", f"{metrics['all']['pearson']:.3f}", f"{metrics['all']['dist_mae']:.4f}"],
    ]
    add_table(slide, 0.75, 1.45, 11.8, 1.5, ["Split", "N", "Ea MAE", "Ea RMSE", "Ea R2", "Pearson", "Dist MAE A"], rows, 10.5)
    add_bullets(
        slide,
        [
            f"Validation Ea MAE is {metrics['val']['mae']:.2f} kcal/mol with R2={metrics['val']['r2']:.3f}.",
            f"Validation geometry distance-matrix MAE is {metrics['val']['dist_mae']:.4f} A.",
            f"Train-to-validation gap is visible: train Ea MAE {metrics['train']['mae']:.2f} vs validation {metrics['val']['mae']:.2f} kcal/mol.",
            f"True Ea range across evaluated data: {metrics['energy_range']['true_min']:.2f} to {metrics['energy_range']['true_max']:.2f} kcal/mol.",
        ],
        1.0,
        3.55,
        5.5,
        2.55,
        15,
    )
    add_callout(slide, 7.05, 3.8, 4.85, 0.82, "Geometry and Ea should be read together: a low barrier error can still hide a poor geometry, and a good geometry may not guarantee a good barrier.", COLORS["orange"], 12)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 13
    slide = blank_slide(prs)
    add_title(slide, "Activation Energy Parity", "Results")
    slide.shapes.add_picture(str(charts["parity"]), Inches(0.85), Inches(1.35), Inches(5.8), Inches(5.25))
    add_bullets(
        slide,
        [
            "The parity plot compares predicted and true activation energies across all evaluated reactions.",
            f"Validation correlation is {metrics['val']['pearson']:.3f}; validation R2 is {metrics['val']['r2']:.3f}.",
            "Most points track the ideal diagonal, while larger deviations concentrate in harder validation cases and high-barrier regions.",
            "Small negative Ea predictions exist in the saved output, which suggests a final non-negative clamp may be useful for deployment.",
        ],
        7.05,
        1.65,
        5.05,
        4.6,
        15,
    )
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 14
    slide = blank_slide(prs)
    add_title(slide, "Error Distributions", "Results")
    slide.shapes.add_picture(str(charts["errors"]), Inches(0.75), Inches(1.38), Inches(5.85), Inches(3.8))
    slide.shapes.add_picture(str(charts["distance"]), Inches(6.75), Inches(1.38), Inches(5.85), Inches(3.8))
    rows = [
        ["Train Ea p90", f"{metrics['train']['p90_abs_error']:.2f} kcal/mol"],
        ["Val Ea p90", f"{metrics['val']['p90_abs_error']:.2f} kcal/mol"],
        ["Train dist p90", f"{metrics['train']['dist_p90']:.4f} A"],
        ["Val dist p90", f"{metrics['val']['dist_p90']:.4f} A"],
    ]
    add_table(slide, 2.0, 5.78, 9.3, 0.95, ["Statistic", "Value"], rows, 10)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 15
    slide = blank_slide(prs)
    add_title(slide, "Neural Ea Head vs Physics Baseline", "Results")
    slide.shapes.add_picture(str(charts["baseline"]), Inches(0.85), Inches(1.45), Inches(6.2), Inches(4.05))
    rows = [
        ["Train", f"{metrics['train']['mae']:.2f}", f"{metrics['physics_train']['mae']:.2f}", f"{metrics['train']['r2']:.3f}", f"{metrics['physics_train']['r2']:.3f}"],
        ["Validation", f"{metrics['val']['mae']:.2f}", f"{metrics['physics_val']['mae']:.2f}", f"{metrics['val']['r2']:.3f}", f"{metrics['physics_val']['r2']:.3f}"],
        ["All", f"{metrics['all']['mae']:.2f}", f"{metrics['physics_all']['mae']:.2f}", f"{metrics['all']['r2']:.3f}", f"{metrics['physics_all']['r2']:.3f}"],
    ]
    add_table(slide, 7.35, 1.55, 4.85, 1.55, ["Split", "Neural MAE", "Physics MAE", "Neural R2", "Physics R2"], rows, 9.5)
    add_bullets(
        slide,
        [
            "The neural Ea head substantially outperforms the simple physics baseline on all splits.",
            "The baseline remains valuable as a sanity check and a lower-complexity comparator.",
            "The validation gap suggests further regularization, split analysis, or reaction-family stratification should be considered.",
        ],
        7.45,
        3.65,
        4.45,
        2.35,
        13.5,
    )
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 16
    slide = blank_slide(prs)
    add_title(slide, "Prediction and Inference Workflow", "Usage")
    add_bullets(
        slide,
        [
            "The predict command loads reactant and product Gaussian log files.",
            "It verifies equal atom counts and consistent atom ordering/types.",
            "It builds the same normalized tensors used in training.",
            "The checkpoint predicts a TS distance matrix, recovers 3D coordinates, predicts Ea, and writes JSON/XYZ outputs.",
            "The saved XYZ file can be used as an initial TS guess in external quantum-chemistry workflows.",
        ],
        0.85,
        1.55,
        6.4,
        4.6,
        15.5,
    )
    code = (
        "python psi_cloud_pipeline.py predict \\\n"
        "  -r reactant.log \\\n"
        "  -p product.log \\\n"
        "  -o prediction.json \\\n"
        "  --xyz predicted_ts.xyz"
    )
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.65), Inches(1.85), Inches(4.35), Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(28, 36, 49)
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["white"]
    add_flow_box(slide, 7.65, 4.05, 4.35, 0.9, "Prediction outputs", "Ea_pred, Ea_pred_physics, D_pred, validation issues, predicted TS coordinates", COLORS["teal"])
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 17
    slide = blank_slide(prs)
    add_title(slide, "Dashboard and Reporting", "Visualization")
    add_bullets(
        slide,
        [
            "The dashboard is generated from detailed_analysis.json and saved as psi_results_dashboard.html.",
            "It includes summary cards, activation-energy parity, distance-MAE histograms, worst-case tables, and representative structures.",
            "3Dmol.js overlays true, predicted, and midpoint/guess geometries for qualitative inspection.",
            "The HTML artifact makes the model reviewable without rerunning training.",
        ],
        0.85,
        1.55,
        6.35,
        4.8,
        16,
    )
    add_flow_box(slide, 7.55, 1.75, 4.35, 0.88, "Quantitative view", "MAE, RMSE, R2, Pearson, split-wise tables", COLORS["blue"])
    add_flow_box(slide, 7.55, 2.9, 4.35, 0.88, "Diagnostic view", "Worst cases, error distributions, geometry histograms", COLORS["orange"])
    add_flow_box(slide, 7.55, 4.05, 4.35, 0.88, "Structural view", "Interactive molecular overlays for selected reactions", COLORS["teal"])
    add_callout(slide, 7.55, 5.35, 4.35, 0.66, "The PPT summarizes the same underlying artifacts but is easier to present in a project review.", COLORS["purple"], 11.5)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 18
    slide = blank_slide(prs)
    add_title(slide, "Strengths, Risks, and Limitations", "Assessment")
    rows = [
        ["Strength", "End-to-end local pipeline with train, predict, and dashboard commands"],
        ["Strength", "Distance-matrix target gives rotational/translation invariance"],
        ["Strength", "EGNN stage adds equivariant coordinate refinement"],
        ["Risk", "Validation error is much higher than training error, so generalization needs scrutiny"],
        ["Risk", "Inference requires consistent atom ordering between reactant and product"],
        ["Risk", "Post-processing consistency is listed in PLANNING.md as a correctness priority"],
        ["Limitation", "Current evaluated atom vocabulary is C/H/N/O only"],
        ["Limitation", "Ea definition uses the higher-energy endpoint, not always the forward barrier"],
    ]
    add_table(slide, 0.8, 1.42, 11.8, 4.75, ["Type", "Item"], rows, 11)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 19
    slide = blank_slide(prs)
    add_title(slide, "Improvement Roadmap", "Next Steps")
    add_flow_box(slide, 0.85, 1.55, 3.45, 1.0, "1. Correctness", "Unify train/eval/inference post-processing; clarify forward-vs-reverse Ea definition.", COLORS["red"])
    add_flow_box(slide, 4.95, 1.55, 3.45, 1.0, "2. Reproducibility", "Seed augmentation, record package versions, keep config snapshots with every run.", COLORS["orange"])
    add_flow_box(slide, 9.05, 1.55, 3.45, 1.0, "3. Model quality", "Stratify validation by reaction family, tune regularization, evaluate uncertainty.", COLORS["blue"])
    add_flow_box(slide, 0.85, 3.35, 3.45, 1.0, "4. Data scaling", "Train on more complete triplets and extend beyond C/H/N/O chemistry.", COLORS["teal"])
    add_flow_box(slide, 4.95, 3.35, 3.45, 1.0, "5. Deployment", "Add non-negative Ea clamp, input validation reports, batch prediction mode.", COLORS["green"])
    add_flow_box(slide, 9.05, 3.35, 3.45, 1.0, "6. Reporting", "Export dashboard figures and representative XYZ overlays for papers/reviews.", COLORS["purple"])
    add_callout(slide, 1.3, 5.55, 10.6, 0.66, "Highest priority: remove train/inference drift and make the activation-energy target definition explicit before relying on new benchmarks.", COLORS["red"], 12.5)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 20
    slide = blank_slide(prs)
    add_title(slide, "Conclusion", "Summary")
    add_bullets(
        slide,
        [
            "The project implements a complete PSI pipeline for transition-state geometry and activation-energy prediction.",
            f"The current evaluated run covers {metrics['all']['count']:,} reactions with a {metrics['train']['count']:,}/{metrics['val']['count']:,} train/validation split.",
            f"Validation performance is {metrics['val']['mae']:.2f} kcal/mol Ea MAE, R2={metrics['val']['r2']:.3f}, and {metrics['val']['dist_mae']:.4f} A geometry distance MAE.",
            "The neural Ea head improves substantially over the physics baseline but still shows a train/validation gap.",
            "The next engineering focus should be correctness, reproducibility, and generalization analysis.",
        ],
        1.0,
        1.55,
        10.8,
        4.5,
        18,
    )
    add_callout(slide, 1.25, 6.0, 10.5, 0.62, "Bottom line: the pipeline is presentation-ready as a research prototype, with clear next steps before production or publication-quality claims.", COLORS["teal"], 13)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 21
    slide = blank_slide(prs)
    add_title(slide, "Appendix: Geometry Preprocessing Details", "Technical Method")
    add_bullets(
        slide,
        [
            "Each reaction is converted into matched R/P/TS triplets with a fixed atom order; samples above max_atoms are skipped.",
            "Connectivity is inferred from coordinates using covalent-radius cutoffs: d_ij <= fragment_bond_scale * (r_i + r_j).",
            "Connected components become molecular fragments; the side with more fragments is used so dissociation/association cases keep separate rigid groups.",
            "For every fragment, reactant coordinates are rigidly aligned onto product coordinates before creating the midpoint geometry.",
            "The midpoint input is c_I = 0.5 * (Kabsch(c_R fragment, c_P fragment) + c_P fragment); its distances are fed beside D_R and D_P.",
            "The geometry mask is block-diagonal over fragments, so padding and unrelated inter-fragment distances do not dominate the training target.",
        ],
        0.75,
        1.48,
        6.75,
        4.95,
        13.7,
        gap=3,
    )
    rows = [
        ["fragment_bond_scale", metric_cfg(metrics, "fragment_bond_scale", 1.45), "Bond cutoff multiplier"],
        ["max_atoms", metric_cfg(metrics, "max_atoms", 30), "Fixed padded tensor size"],
        ["spectator_threshold", f"{metric_cfg(metrics, 'spectator_threshold', 0.15)} A", "Near-unchanged pair cutoff"],
        ["geom_mask", "fragment blocks", "Supervised pairwise geometry region"],
    ]
    add_table(slide, 7.85, 1.55, 4.65, 1.95, ["Object", "Value", "Role"], rows, 9.3)
    add_flow_box(slide, 7.85, 3.92, 4.65, 0.84, "Why align first?", "It removes arbitrary global rotation/translation before constructing the midpoint guess.", COLORS["blue"])
    add_flow_box(slide, 7.85, 5.0, 4.65, 0.84, "Why fragment-aware?", "Separate fragments should not be forced into one artificial rigid body during alignment or MDS.", COLORS["teal"])
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 22
    slide = blank_slide(prs)
    add_title(slide, "Appendix: Kabsch Alignment", "Technical Method")
    add_bullets(
        slide,
        [
            "For a fragment with at least two atoms, both coordinate sets are centered by subtracting their fragment centroids.",
            "The cross-covariance matrix C = P_centered^T Q_centered is decomposed with SVD: C = V S W.",
            "The rotation is R = V W, except when det(V W) < 0; then the last column of V is sign-flipped to avoid a reflected geometry.",
            "Aligned reactant coordinates are returned as P_centered R + mean(Q), so the fragment lands in the product frame.",
            "Single-atom fragments are copied directly to the product coordinate because rotation is undefined for one point.",
        ],
        0.78,
        1.48,
        6.25,
        4.58,
        14.2,
        gap=4,
    )
    add_code_box(
        slide,
        7.35,
        1.58,
        4.95,
        2.95,
        "P0 = P - mean(P)\n"
        "Q0 = Q - mean(Q)\n"
        "C  = P0.T @ Q0\n"
        "V, S, W = svd(C)\n"
        "if det(V @ W) < 0:\n"
        "    V[:, -1] *= -1\n"
        "R = V @ W\n"
        "P_aligned = P0 @ R + mean(Q)",
        11.2,
    )
    add_callout(
        slide,
        7.35,
        5.05,
        4.95,
        0.88,
        "This is a rigid-body operation: bond lengths inside the reactant fragment are preserved; only orientation and translation change.",
        COLORS["teal"],
        11.6,
    )
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 23
    slide = blank_slide(prs)
    add_title(slide, "Appendix: Distance Head and Clamping", "Technical Method")
    add_bullets(
        slide,
        [
            "The geometry head builds pair features from transformer node context, atom embeddings, physical descriptors, and raw R/I/P distances.",
            "It predicts two values per atom pair: alpha_logit and delta; alpha = sigmoid(alpha_logit) stays between 0 and 1.",
            "The base distance is an endpoint interpolation: D_base = alpha * D_R + (1 - alpha) * D_P.",
            f"The learned correction is clipped to +/-{metric_cfg(metrics, 'delta_clamp', 3.0)} A before addition, preventing one pair from making an unphysical jump.",
            "After addition, distances are clamped to non-negative values, symmetrized, diagonal-zeroed, and multiplied by the valid atom mask.",
            "At inference, steric collision clamping enforces D_ij >= 0.75 * (r_i + r_j) for every non-diagonal atom pair.",
        ],
        0.78,
        1.48,
        6.45,
        4.9,
        13.4,
        gap=3,
    )
    add_code_box(
        slide,
        7.55,
        1.55,
        4.72,
        2.55,
        "alpha = sigmoid(out[..., 0])\n"
        "delta = clamp(out[..., 1], -delta_clamp, +delta_clamp)\n"
        "D_base = alpha * D_R + (1 - alpha) * D_P\n"
        "D_TS = clamp(D_base + delta, min=0)\n"
        "D_TS = 0.5 * (D_TS + D_TS.T)\n"
        "diag(D_TS) = 0",
        10.3,
    )
    rows = [
        ["Delta clamp", f"+/-{metric_cfg(metrics, 'delta_clamp', 3.0)} A", "Limits learned correction"],
        ["Steric floor", "0.75*(r_i+r_j)", "Prevents atom overlap"],
        ["Symmetry", "0.5*(D+D^T)", "Valid distance matrix form"],
        ["Diagonal", "0", "Self-distance constraint"],
    ]
    add_table(slide, 7.55, 4.55, 4.72, 1.48, ["Check", "Rule", "Purpose"], rows, 8.7)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 24
    slide = blank_slide(prs)
    add_title(slide, "Appendix: Spectator Bounds and Triangle Inequality", "Technical Method")
    add_bullets(
        slide,
        [
            "Spectator pairs are atom pairs whose endpoint distance barely changes: abs(D_R - D_P) < spectator_threshold.",
            f"During inference, masked spectator pairs are clamped within +/-{metric_cfg(metrics, 'spectator_tol', 0.05) * 100:.0f}% of their endpoint midpoint distance.",
            "During training, spectator pairs add a matrix-wise physics loss that pulls predicted TS distances toward D_I.",
            "Active-pair bounds penalize distances outside [min(D_R,D_P)-0.2, max(D_R,D_P)+0.2] using squared ReLU excess.",
            "Triangle inequality is enforced fragment-by-fragment: if D_ij > D_ik + D_kj + tol, the direct distance is shortened to the two-edge shortcut.",
            "The fragment scope matters because disconnected species may have arbitrary relative placement; enforcing all cross-fragment triangles would inject false geometry.",
        ],
        0.75,
        1.45,
        6.85,
        5.12,
        12.6,
        gap=2.4,
    )
    add_code_box(
        slide,
        7.95,
        1.52,
        4.48,
        2.7,
        "spectator = abs(D_R - D_P) < threshold\n"
        "d_ref = 0.5 * (D_R + D_P)\n"
        "D_ij = clip(D_ij,\n"
        "            d_ref * (1 - tol),\n"
        "            d_ref * (1 + tol))\n\n"
        "if D[i,j] > D[i,k] + D[k,j] + tri_tol:\n"
        "    D[i,j] = D[j,i] = D[i,k] + D[k,j]",
        9.7,
    )
    rows = [
        ["Spectator threshold", f"{metric_cfg(metrics, 'spectator_threshold', 0.15)} A", "Detect unchanged pairs"],
        ["Spectator tolerance", f"{metric_cfg(metrics, 'spectator_tol', 0.05) * 100:.0f}%", "Inference clamp band"],
        ["Triangle tolerance", "0.05 A", "Ignore tiny numerical violations"],
        ["Active bounds margin", "0.2 A", "Soft training penalty"],
    ]
    add_table(slide, 7.95, 4.72, 4.48, 1.5, ["Term", "Value", "Use"], rows, 8.7)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 25
    slide = blank_slide(prs)
    add_title(slide, "Appendix: Distance-to-Coordinate Recovery", "Technical Method")
    add_bullets(
        slide,
        [
            "The coarse TS distance matrix is embedded into 3D with classical MDS, using masked double-centering so padded atoms do not affect the embedding.",
            "Eigenvalues are sorted, the top three non-negative components are kept, and coordinates are scaled by sqrt(eigenvalue).",
            "For disconnected systems, MDS is also applied by fragment during post-processing, then each fragment is Kabsch-aligned onto the midpoint reference.",
            "The EGNN starts from the differentiable MDS coordinates and refines them with equivariant message passing.",
            f"Each EGNN layer clips coordinate displacement messages to +/-{metric_cfg(metrics, 'egnn_coord_clamp', 2.0)} A, then averages by node degree.",
            "The final pairwise distances are recomputed from refined coordinates and used for the main geometry loss and downstream Ea head.",
        ],
        0.75,
        1.5,
        6.85,
        4.8,
        13.2,
        gap=3,
    )
    add_flow_box(slide, 8.0, 1.62, 4.35, 0.76, "1. Predicted D_TS", "Symmetric, non-negative pairwise distances", COLORS["blue"])
    add_flow_box(slide, 8.0, 2.55, 4.35, 0.76, "2. MDS embedding", "Double-center D^2 and keep top 3 eigenvectors", COLORS["teal"])
    add_flow_box(slide, 8.0, 3.48, 4.35, 0.76, "3. EGNN refinement", "Equivariant coordinate updates with clipped messages", COLORS["purple"])
    add_flow_box(slide, 8.0, 4.41, 4.35, 0.76, "4. Final distances", "Recompute ||x_i - x_j|| for loss and reporting", COLORS["orange"])
    add_callout(slide, 8.0, 5.55, 4.35, 0.72, "The model learns distances first because distances are invariant; coordinates are recovered only after the pairwise geometry is stable.", COLORS["teal"], 11.4)
    add_footer(slide, slide_no)
    slide_no += 1

    # Slide 26
    slide = blank_slide(prs)
    add_title(slide, "Appendix: Exact Training Objective", "Technical Method")
    add_bullets(
        slide,
        [
            "The main geometry term is Huber loss(delta=0.5) between EGNN-refined predicted distances and true TS distances over valid non-diagonal pairs.",
            f"The coarse auxiliary term applies the same Huber loss to the pre-EGNN distance head output with weight {metric_cfg(metrics, 'geom_coarse_weight', 0.5)}.",
            "The spectator physics term is MSE between predicted TS distances and midpoint distances for pairs with abs(D_R - D_P) below threshold.",
            "The active-bounds term penalizes only violations beyond the endpoint interval plus a 0.2 A margin.",
            "Total geometry loss is: L_geom + geom_coarse_weight * L_coarse + 0.2 * (L_spectator + 0.5 * L_bounds).",
            f"The Ea loss is SmoothL1 on normalized Ea and is added only after epoch {metric_cfg(metrics, 'ea_warmup_epochs', 200)}.",
            f"Gradients are clipped to norm {metric_cfg(metrics, 'grad_clip', 1.0)} before the optimizer step.",
        ],
        0.75,
        1.4,
        6.95,
        5.45,
        12.4,
        gap=2,
    )
    add_code_box(
        slide,
        7.95,
        1.48,
        4.5,
        2.75,
        "L_main   = Huber(D_refined, D_TS)\n"
        "L_coarse = Huber(D_coarse, D_TS)\n"
        "L_spec   = MSE(D_refined * M_spec, D_I * M_spec)\n"
        "L_bounds = sum(ReLU(D-minmax_violation)^2)\n\n"
        "L_pinn = L_spec + 0.5 * L_bounds\n"
        "L = L_main + w_coarse*L_coarse + 0.2*L_pinn\n"
        "if epoch > ea_warmup: L += w_ea * SmoothL1(Ea_norm)",
        9.2,
    )
    rows = [
        ["Optimizer", "AdamW", "Cosine schedule with warmup"],
        ["AMP", metric_cfg(metrics, "amp", True), "Mixed precision when available"],
        ["Geometry mask", "valid non-diagonal pairs", "Padding/self-distance excluded"],
        ["Ea target", "z-scored kcal/mol", "Denormalized for reporting"],
    ]
    add_table(slide, 7.95, 4.75, 4.5, 1.48, ["Item", "Setting", "Why"], rows, 8.5)
    add_footer(slide, slide_no)


def add_algorithm_focused_deck_slides(prs, metrics, charts):
    slide_no = 1
    cfg = metrics["config"]
    raw = metrics["raw_summary"]

    def finish(slide):
        nonlocal slide_no
        add_footer(slide, slide_no)
        slide_no += 1

    # Slide 1
    slide = blank_slide(prs)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(236, 244, 250)
    shape.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.23))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["teal"]
    band.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.8), Inches(0.95))
    p = box.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = "Aptos Display"
    p.font.size = Pt(39)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    sub = slide.shapes.add_textbox(Inches(0.85), Inches(2.05), Inches(11.1), Inches(0.9))
    p = sub.text_frame.paragraphs[0]
    p.text = "Algorithm-focused implementation review: data triplets, invariant geometry learning, EGNN coordinate refinement, Ea regression, and inference constraints."
    p.font.name = "Aptos"
    p.font.size = Pt(17)
    p.font.color.rgb = COLORS["muted"]
    add_number_card(slide, 0.85, 3.75, 2.15, 0.95, fmt_int(metrics["all"]["count"]), "evaluated reactions", COLORS["blue"])
    add_number_card(slide, 3.25, 3.75, 2.15, 0.95, f"{metrics['val']['mae']:.2f}", "val Ea MAE kcal/mol", COLORS["teal"])
    add_number_card(slide, 5.65, 3.75, 2.15, 0.95, f"{metrics['val']['dist_mae']:.4f}", "val distance MAE A", COLORS["orange"])
    add_number_card(slide, 8.05, 3.75, 2.15, 0.95, f"{metrics['val']['r2']:.3f}", "val Ea R2", COLORS["purple"])
    add_callout(
        slide,
        0.85,
        5.35,
        11.1,
        0.78,
        "Core idea: learn the TS distance matrix from reactant/product distances, embed it into 3D, refine coordinates equivariantly, then predict Ea from the refined TS representation.",
        COLORS["blue"],
        13.2,
    )
    finish(slide)

    # Slide 2
    slide = blank_slide(prs)
    add_title(slide, "Data Assembly and Targets", "Implementation")
    add_bullets(
        slide,
        [
            "build_reaction_samples reads extracted_dataset.json once, groups filenames by reaction ID, and keeps only complete reactant/product/TS triplets.",
            "Atom order is taken from the TS entry and reused for reactant and product tensors; samples above max_atoms are skipped before padding.",
            "Ea target is (E_TS - max(E_R, E_P)) * 627.509 kcal/mol. Negative targets are skipped when skip_negative_ea is enabled.",
            "The true TS target is a padded D_TS pairwise distance matrix, so supervision is invariant to global rotation and translation.",
            "A TS-fragment geometry mask is precomputed from covalent-radius connectivity so loss ignores padding and disconnected cross-fragment distances.",
        ],
        0.75,
        1.45,
        6.8,
        4.9,
        13.6,
        gap=3,
    )
    add_number_card(slide, 8.0, 1.5, 1.95, 0.78, fmt_int(raw.get("raw_entries", 0)), "raw entries", COLORS["blue"])
    add_number_card(slide, 10.2, 1.5, 1.95, 0.78, fmt_int(raw.get("complete_triplets", 0)), "complete triplets", COLORS["green"])
    rows = [
        ["Train", fmt_int(metrics["train"]["count"]), f"{metrics['train']['atoms_mean']:.1f}", f"{metrics['train']['atoms_min']}-{metrics['train']['atoms_max']}"],
        ["Validation", fmt_int(metrics["val"]["count"]), f"{metrics['val']['atoms_mean']:.1f}", f"{metrics['val']['atoms_min']}-{metrics['val']['atoms_max']}"],
        ["All", fmt_int(metrics["all"]["count"]), f"{metrics['all']['atoms_mean']:.1f}", f"{metrics['all']['atoms_min']}-{metrics['all']['atoms_max']}"],
    ]
    add_table(slide, 7.85, 2.75, 4.35, 1.35, ["Split", "N", "Avg atoms", "Range"], rows, 9.6)
    add_callout(slide, 7.85, 4.75, 4.35, 0.78, "This slide replaces broad motivation with the actual target definition and sample filtering logic used by the code.", COLORS["teal"], 11.2)
    finish(slide)

    # Slide 3
    slide = blank_slide(prs)
    add_title(slide, "Feature Tensor Construction", "Implementation")
    add_bullets(
        slide,
        [
            "ReactionDataset optionally adds Gaussian coordinate noise only to R/P coordinates during training; validation uses clean coordinates.",
            "D_R, D_P, and D_I are built from pairwise Euclidean distances. D_I is the midpoint distance matrix used as a stable geometric prior.",
            "atom_ids feed a learned embedding; atom_phys adds normalized electronegativity, atomic number, and atomic mass.",
            "energy_feats is a normalized 20D vector: composition, signed reaction energy, displacement statistics, and bond-angle statistics.",
            "Normalization statistics are computed only on training indices to prevent validation leakage.",
        ],
        0.75,
        1.5,
        6.65,
        4.65,
        14,
        gap=3,
    )
    rows = [
        ["D_R", "[B,30,30]", "Reactant distances"],
        ["D_I", "[B,30,30]", "Midpoint/interpolation prior"],
        ["D_P", "[B,30,30]", "Product distances"],
        ["mask", "[B,30]", "Real atoms vs padding"],
        ["geom_mask", "[B,30,30]", "Pairs included in geometry loss"],
        ["energy_feats", "[B,20]", "Global chemistry descriptor"],
    ]
    add_table(slide, 7.75, 1.55, 4.65, 2.65, ["Tensor", "Shape", "Role"], rows, 8.7)
    add_code_box(
        slide,
        7.75,
        4.65,
        4.65,
        1.45,
        "D_R = pairwise(c_R)\n"
        "D_P = pairwise(c_P)\n"
        "D_I = 0.5 * (D_R + D_P)\n"
        "Ea_norm = (Ea - mean_train) / std_train",
        10.3,
    )
    finish(slide)

    # Slide 4
    slide = blank_slide(prs)
    add_title(slide, "End-to-End Algorithm", "Workflow")
    x0, y0, w, h, gap = 0.58, 1.62, 1.78, 1.0, 0.19
    boxes = [
        ("1. Logs", "Gaussian R/P/TS energies and coordinates", COLORS["blue"]),
        ("2. Triplets", "Complete reaction samples with fixed atom order", COLORS["teal"]),
        ("3. Tensors", "D_R, D_I, D_P, masks, atom and energy features", COLORS["green"]),
        ("4. PSI core", "Temporal R/I/P encoding plus atom attention", COLORS["purple"]),
        ("5. Geometry", "Alpha interpolation and clamped pair delta", COLORS["orange"]),
        ("6. 3D/Ea", "MDS, EGNN refinement, learned barrier head", COLORS["red"]),
    ]
    for i, (title, body, color) in enumerate(boxes):
        add_flow_box(slide, x0 + i * (w + gap), y0, w, h, title, body, color)
        if i < len(boxes) - 1:
            add_arrow(slide, x0 + (i + 1) * w + i * gap + 0.02, y0 + 0.38, 0.2)
    add_bullets(
        slide,
        [
            "Training path: predict D_TS, embed to coordinates, refine with EGNN, compute geometry losses, and add Ea loss after warmup.",
            "Evaluation path: post-process the predicted distance matrix, recover fragment-aware coordinates, score geometry and Ea, and store detailed_analysis.json.",
            "Inference path: reuse the same normalization metadata from the checkpoint, then write JSON plus optional XYZ for quantum-chemistry follow-up.",
        ],
        1.0,
        3.35,
        10.8,
        2.05,
        14.2,
        gap=4,
    )
    add_callout(slide, 1.05, 5.9, 10.8, 0.64, "The model is not a direct coordinate regressor: it learns invariant distances first, then converts to coordinates only after the pairwise geometry is stable.", COLORS["teal"], 12.3)
    finish(slide)

    # Slide 5
    slide = blank_slide(prs)
    add_title(slide, "PSI Core Encoder", "Architecture")
    add_bullets(
        slide,
        [
            "Each atom row sees all pairwise distances to the other 30 padded positions through 32 Gaussian radial basis channels.",
            "For every atom, the model builds three state embeddings: reactant, midpoint, and product. Each state also includes atom embedding and physical descriptors.",
            "A 2-layer bidirectional GRU runs over the R -> I -> P sequence per atom; the midpoint hidden state becomes the atom context.",
            "Three pre-norm Transformer encoder layers then let atoms exchange global context through self-attention while honoring the padding mask.",
            "The output is a per-atom latent vector used by the pairwise geometry head.",
        ],
        0.75,
        1.45,
        6.7,
        4.85,
        13.4,
        gap=3,
    )
    add_flow_box(slide, 7.85, 1.55, 4.25, 0.78, "GaussianEmbedding", f"{cfg.get('n_gaussians', 32)} radial basis channels from {cfg.get('gauss_start', 0.4)}-{cfg.get('gauss_stop', 6.0)} A", COLORS["blue"])
    add_flow_box(slide, 7.85, 2.55, 4.25, 0.78, "BiGRU", f"{cfg.get('gru_layers', 2)} layers, hidden={cfg.get('gru_hidden', 128)}, reads R/I/P as a short geometry trajectory", COLORS["teal"])
    add_flow_box(slide, 7.85, 3.55, 4.25, 0.78, "Transformer", f"{cfg.get('attn_layers', 3)} layers, {cfg.get('attn_heads', 8)} heads, ff_dim={cfg.get('ff_dim', 512)}", COLORS["purple"])
    add_code_box(
        slide,
        7.85,
        4.75,
        4.25,
        1.15,
        "state_seq = [emb_R, emb_I, emb_P]\n"
        "context = BiGRU(state_seq)[:, midpoint]\n"
        "node_feats = Transformer(context, pad_mask)",
        9.5,
    )
    finish(slide)

    # Slide 6
    slide = blank_slide(prs)
    add_title(slide, "Pairwise Geometry Head", "Architecture")
    add_bullets(
        slide,
        [
            "For each atom pair (i,j), the head concatenates node context for i and j, atom features for i and j, and raw distances D_R, D_I, D_P.",
            "The MLP predicts alpha_logit and delta. alpha = sigmoid(alpha_logit) forces interpolation weight into [0,1].",
            "The initial final layer is zeroed, so the first prediction starts at the endpoint midpoint with no learned correction.",
            "delta is clipped before addition, which prevents a single pairwise prediction from creating an extreme unphysical distance.",
            "The matrix is symmetrized, diagonal-zeroed, non-negative clamped, and multiplied by the valid atom mask.",
        ],
        0.75,
        1.45,
        6.65,
        4.9,
        13.4,
        gap=3,
    )
    add_section_label(slide, "Distance formula", 7.75, 1.55, 4.5)
    add_code_box(
        slide,
        7.75,
        1.95,
        4.65,
        2.0,
        "alpha = sigmoid(alpha_logit)\n"
        f"delta = clamp(delta, -{metric_cfg(metrics, 'delta_clamp', 3.0)}, +{metric_cfg(metrics, 'delta_clamp', 3.0)})\n"
        "D_base = alpha*D_R + (1-alpha)*D_P\n"
        "D_TS = symmetrize(clamp(D_base + delta, 0))",
        10.0,
    )
    rows = [
        ["Input", "node_i, node_j, atom_i, atom_j, D_R, D_I, D_P"],
        ["Output", "coarse TS distance matrix"],
        ["Aux loss", "Huber(coarse D_TS, true D_TS)"],
    ]
    add_table(slide, 7.75, 4.55, 4.65, 1.35, ["Part", "Implementation"], rows, 9.0)
    finish(slide)

    # Slide 7
    slide = blank_slide(prs)
    add_title(slide, "MDS and EGNN Refinement", "Architecture")
    add_bullets(
        slide,
        [
            "torch_mds_coords turns the coarse distance matrix into 3D coordinates using masked classical MDS and the top three non-negative eigen-components.",
            "The EGNN consumes atom chemical features plus the initial MDS coordinates. Messages depend on squared pairwise distances, so feature updates remain invariant.",
            "Coordinate updates are equivariant: if the input geometry is rotated or translated, the output rotates or translates consistently.",
            "Each EGNN layer clips coordinate displacement messages and normalizes by node degree, preventing unstable geometry jumps.",
            "The final training distance matrix is recomputed from refined coordinates, so the main geometry loss supervises the actual 3D refinement.",
        ],
        0.75,
        1.45,
        6.85,
        4.95,
        13.2,
        gap=3,
    )
    add_flow_box(slide, 8.0, 1.55, 4.25, 0.72, "1. Coarse D_TS", "Pairwise output from geometry head", COLORS["blue"])
    add_flow_box(slide, 8.0, 2.45, 4.25, 0.72, "2. Classical MDS", "Double-center D^2 and keep top 3 components", COLORS["teal"])
    add_flow_box(slide, 8.0, 3.35, 4.25, 0.72, "3. EGNN layers", f"{cfg.get('egnn_layers', 4)} layers, hidden={cfg.get('egnn_hidden', 128)}, coord clamp={metric_cfg(metrics, 'egnn_coord_clamp', 2.0)} A", COLORS["purple"])
    add_flow_box(slide, 8.0, 4.25, 4.25, 0.72, "4. Refined D_TS", "Recomputed from x_TS for loss, metrics, and Ea", COLORS["orange"])
    add_callout(slide, 8.0, 5.35, 4.25, 0.74, "This is the technical reason the pipeline can learn geometry without depending on an arbitrary Cartesian frame.", COLORS["teal"], 11.2)
    finish(slide)

    # Slide 8
    slide = blank_slide(prs)
    add_title(slide, "Activation Energy Modeling", "Architecture")
    add_bullets(
        slide,
        [
            "EaHead uses EGNN-refined per-atom features, masked mean pooling, z-scored signed reaction energy, and normalized 20D energy features.",
            "The head predicts normalized Ea. Reporting denormalizes with ea_mean and ea_std saved inside the checkpoint metadata.",
            "Ea loss starts after ea_warmup_epochs so early poor geometries do not dominate the shared EGNN representation.",
            "A Marcus/Hammond/OLS physics baseline is fit after geometry prediction and saved for comparison, not used as the primary prediction when the neural head exists.",
            "The baseline features include geometric reorganization energy, Hammond index, signed reaction energy, and a fitted intercept.",
        ],
        0.75,
        1.45,
        6.9,
        4.9,
        13.3,
        gap=3,
    )
    rows = [
        ["Neural Ea head", f"Val MAE {metrics['val']['mae']:.2f}", "Primary output"],
        ["Physics baseline", f"Val MAE {metrics['physics_val']['mae']:.2f}", "Comparator/sanity check"],
        ["Warmup", f"{metric_cfg(metrics, 'ea_warmup_epochs', 200)} epochs", "Geometry first"],
        ["Selection weight", metric_cfg(metrics, "ea_select_weight", 0.25), "Blends Ea into checkpoint score"],
    ]
    add_table(slide, 8.0, 1.65, 4.25, 1.95, ["Component", "Value", "Role"], rows, 8.8)
    add_code_box(
        slide,
        8.0,
        4.15,
        4.25,
        1.45,
        "mol = masked_mean(h_ts)\n"
        "feat = concat(mol, de_rxn_norm, energy_feats)\n"
        "Ea = ea_head(feat) * ea_std + ea_mean",
        9.8,
    )
    finish(slide)

    # Slide 9
    slide = blank_slide(prs)
    add_title(slide, "Training Objective", "Training")
    add_bullets(
        slide,
        [
            "Main geometry loss: Huber(delta=0.5) between EGNN-refined predicted distances and true TS distances over valid non-diagonal atom pairs.",
            "Coarse auxiliary loss: same Huber loss on the pre-EGNN distance head, weighted by geom_coarse_weight, so the MDS seed improves directly.",
            "Spectator physics term: pairs with abs(D_R - D_P) below the threshold are pulled toward D_I because those distances should remain nearly unchanged.",
            "Active-pair bounds term: penalizes predicted distances outside the endpoint interval plus a 0.2 A margin using squared ReLU violations.",
            "Ea term: SmoothL1 on normalized Ea, added after warmup, with gradient flowing through the shared EGNN features.",
        ],
        0.72,
        1.42,
        6.95,
        5.0,
        12.8,
        gap=2.5,
    )
    add_code_box(
        slide,
        7.95,
        1.5,
        4.65,
        2.75,
        "L_main   = Huber(D_refined, D_TS)\n"
        "L_coarse = Huber(D_coarse, D_TS)\n"
        "L_spec   = MSE(D_refined, D_I) on spectator pairs\n"
        "L_bounds = ReLU(endpoint interval violation)^2\n\n"
        "L = L_main + w_coarse*L_coarse + 0.2*(L_spec + 0.5*L_bounds)\n"
        "if epoch > ea_warmup: L += w_ea*SmoothL1(Ea_norm)",
        8.9,
    )
    rows = [
        ["geom_coarse_weight", metric_cfg(metrics, "geom_coarse_weight", 0.5)],
        ["spectator_threshold", f"{metric_cfg(metrics, 'spectator_threshold', 0.15)} A"],
        ["ea_loss_weight", metric_cfg(metrics, "ea_loss_weight", 1.0)],
        ["grad_clip", metric_cfg(metrics, "grad_clip", 1.0)],
    ]
    add_table(slide, 7.95, 4.75, 4.65, 1.38, ["Parameter", "Value"], rows, 9.0)
    finish(slide)

    # Slide 10
    slide = blank_slide(prs)
    add_title(slide, "Training Configuration and Dynamics", "Training")
    slide.shapes.add_picture(str(charts["training"]), Inches(0.75), Inches(1.35), Inches(6.4), Inches(3.85))
    rows = [
        ["Target reactions", fmt_int(cfg.get("target_reactions", metrics["all"]["count"]))],
        ["Train / validation", f"{metrics['train']['count']:,} / {metrics['val']['count']:,}"],
        ["Batch size", cfg.get("batch_size", 32)],
        ["Learning rate", cfg.get("lr", 1.5e-4)],
        ["Weight decay", cfg.get("weight_decay", 1e-2)],
        ["Epoch cap", cfg.get("epochs", 1500)],
        ["Patience", cfg.get("patience", 120)],
        ["AMP", cfg.get("amp", True)],
    ]
    add_table(slide, 7.55, 1.35, 4.55, 2.9, ["Setting", "Value"], rows, 9.0)
    add_bullets(
        slide,
        [
            f"Best validation geometry: epoch {metrics['best_val_geom']['epoch']} with loss {metrics['best_val_geom']['val_geom']:.5f}.",
            f"Best validation Ea: epoch {metrics['best_val_ea']['epoch']} with MAE {metrics['best_val_ea']['val_ea_mae']:.2f} kcal/mol.",
            f"Final recorded validation geometry loss: {metrics['history'][-1]['val_geom']:.5f}; final validation Ea MAE: {metrics['history'][-1]['val_ea_mae']:.2f} kcal/mol.",
        ],
        7.55,
        4.65,
        4.45,
        1.45,
        11.4,
        gap=2,
    )
    finish(slide)

    # Slide 11
    slide = blank_slide(prs)
    add_title(slide, "Evaluation Readout", "Results")
    rows = [
        ["Train", fmt_int(metrics["train"]["count"]), f"{metrics['train']['mae']:.2f}", f"{metrics['train']['rmse']:.2f}", f"{metrics['train']['r2']:.3f}", f"{metrics['train']['pearson']:.3f}", f"{metrics['train']['dist_mae']:.4f}"],
        ["Validation", fmt_int(metrics["val"]["count"]), f"{metrics['val']['mae']:.2f}", f"{metrics['val']['rmse']:.2f}", f"{metrics['val']['r2']:.3f}", f"{metrics['val']['pearson']:.3f}", f"{metrics['val']['dist_mae']:.4f}"],
        ["All", fmt_int(metrics["all"]["count"]), f"{metrics['all']['mae']:.2f}", f"{metrics['all']['rmse']:.2f}", f"{metrics['all']['r2']:.3f}", f"{metrics['all']['pearson']:.3f}", f"{metrics['all']['dist_mae']:.4f}"],
    ]
    add_table(slide, 0.7, 1.45, 11.95, 1.45, ["Split", "N", "Ea MAE", "Ea RMSE", "Ea R2", "Pearson", "Dist MAE A"], rows, 10.0)
    slide.shapes.add_picture(str(charts["parity"]), Inches(0.75), Inches(3.25), Inches(5.1), Inches(3.0))
    add_bullets(
        slide,
        [
            f"Validation Ea MAE is {metrics['val']['mae']:.2f} kcal/mol with R2={metrics['val']['r2']:.3f}; validation distance MAE is {metrics['val']['dist_mae']:.4f} A.",
            f"Train/validation Ea gap: {metrics['train']['mae']:.2f} vs {metrics['val']['mae']:.2f} kcal/mol, so generalization analysis is still required.",
            "The evaluation file keeps both D_true and D_pred, so geometry errors can be diagnosed pair-by-pair instead of only through scalar summaries.",
        ],
        6.25,
        3.35,
        5.75,
        2.6,
        13.0,
        gap=3,
    )
    finish(slide)

    # Slide 12
    slide = blank_slide(prs)
    add_title(slide, "Error Structure and Baseline", "Results")
    slide.shapes.add_picture(str(charts["errors"]), Inches(0.7), Inches(1.35), Inches(4.0), Inches(2.9))
    slide.shapes.add_picture(str(charts["distance"]), Inches(4.8), Inches(1.35), Inches(4.0), Inches(2.9))
    slide.shapes.add_picture(str(charts["baseline"]), Inches(8.9), Inches(1.35), Inches(3.65), Inches(2.9))
    rows = [
        ["Val Ea p90", f"{metrics['val']['p90_abs_error']:.2f} kcal/mol"],
        ["Val dist p90", f"{metrics['val']['dist_p90']:.4f} A"],
        ["Neural val MAE", f"{metrics['val']['mae']:.2f} kcal/mol"],
        ["Physics val MAE", f"{metrics['physics_val']['mae']:.2f} kcal/mol"],
    ]
    add_table(slide, 1.05, 4.75, 4.3, 1.35, ["Statistic", "Value"], rows, 9.2)
    add_bullets(
        slide,
        [
            "Energy and geometry errors are complementary diagnostics: a barrier can look acceptable while the predicted TS structure is poor.",
            "The physics baseline is intentionally simple; it is useful because it exposes whether the learned head is doing more than fitting reaction energy trends.",
            "The p90 metrics highlight the tail cases that should drive the next round of stratified evaluation.",
        ],
        6.0,
        4.8,
        5.95,
        1.45,
        11.6,
        gap=2,
    )
    finish(slide)

    # Slide 13
    slide = blank_slide(prs)
    add_title(slide, "Inference Path", "Usage")
    add_bullets(
        slide,
        [
            "predict_transition_state loads reactant and product Gaussian logs, then checks atom count and atom type/order consistency.",
            "It rebuilds D_R, D_I, D_P, atom physical features, signed reaction energy, and energy_feats using the checkpoint's saved normalization statistics.",
            "The model outputs a predicted distance matrix and normalized neural Ea; Ea is denormalized from checkpoint metadata.",
            "Post-processing applies spectator clamps, triangle inequality checks, fragment-aware MDS, and optional XYZ export.",
            "The JSON output includes Ea_pred, Ea_pred_physics, Ea_source, D_pred, validation issues, geom_mask, and predicted coordinates.",
        ],
        0.75,
        1.45,
        6.7,
        4.95,
        13.3,
        gap=3,
    )
    add_code_box(
        slide,
        7.85,
        1.7,
        4.55,
        1.5,
        "python psi_cloud_pipeline.py predict \\\n"
        "  -r reactant.log \\\n"
        "  -p product.log \\\n"
        "  -o prediction.json \\\n"
        "  --xyz predicted_ts.xyz",
        10.5,
    )
    add_flow_box(slide, 7.85, 3.85, 4.55, 0.76, "Same feature path", "Inference reuses training-time tensor construction and normalization.", COLORS["teal"])
    add_flow_box(slide, 7.85, 4.85, 4.55, 0.76, "Downstream use", "predicted_ts.xyz is an initial guess, not a replacement for quantum validation.", COLORS["orange"])
    finish(slide)

    # Slide 14
    slide = blank_slide(prs)
    add_title(slide, "Post-Processing Constraints", "Technique")
    add_bullets(
        slide,
        [
            "Spectator pairs are distances with little endpoint change: abs(D_R - D_P) < spectator_threshold. In inference they are clamped near the endpoint midpoint.",
            "Steric floor prevents atom overlap by enforcing a covalent-radius lower bound for non-diagonal pairs.",
            "Triangle inequality is enforced within fragments by shortening impossible direct distances to the best two-edge shortcut.",
            "Fragment-aware MDS reconstructs disconnected systems block-by-block and Kabsch-aligns fragments to the midpoint reference.",
            "validate_ts_geometry reports suspicious active-bond lengths and spectator violations so bad predictions are visible to the caller.",
        ],
        0.75,
        1.42,
        6.85,
        5.05,
        12.9,
        gap=2.5,
    )
    add_code_box(
        slide,
        7.95,
        1.5,
        4.48,
        2.55,
        "spectator = abs(D_R - D_P) < threshold\n"
        "D_ref = 0.5 * (D_R + D_P)\n"
        "D[spectator] = clip(D, D_ref*(1-tol), D_ref*(1+tol))\n\n"
        "if D[i,j] > D[i,k] + D[k,j] + tri_tol:\n"
        "    D[i,j] = D[i,k] + D[k,j]",
        9.4,
    )
    rows = [
        ["Spectator threshold", f"{metric_cfg(metrics, 'spectator_threshold', 0.15)} A"],
        ["Spectator tolerance", f"{metric_cfg(metrics, 'spectator_tol', 0.05) * 100:.0f}%"],
        ["Triangle tolerance", "0.05 A"],
        ["Fragment bond scale", metric_cfg(metrics, "fragment_bond_scale", 1.45)],
    ]
    add_table(slide, 7.95, 4.6, 4.48, 1.45, ["Constraint", "Value"], rows, 9.0)
    finish(slide)

    # Slide 15
    slide = blank_slide(prs)
    add_title(slide, "Kabsch Alignment and Fragments", "Technique")
    add_bullets(
        slide,
        [
            "Connectivity is inferred with covalent radii and fragment_bond_scale. The fragment set is chosen from the side with more fragments so association/dissociation cases remain separated.",
            "For each fragment, reactant coordinates are rigidly aligned onto product coordinates before building midpoint references for energy features and coordinate recovery.",
            "Kabsch alignment centers both point sets, decomposes the covariance matrix with SVD, fixes reflection if det(VW) < 0, and applies the rotation plus product centroid.",
            "Single-atom fragments are copied directly because rotation is undefined for one point.",
            "This preserves internal geometry while removing arbitrary global orientation before midpoint and MDS alignment steps.",
        ],
        0.75,
        1.42,
        6.85,
        5.05,
        12.8,
        gap=2.5,
    )
    add_code_box(
        slide,
        7.95,
        1.55,
        4.55,
        2.6,
        "P0 = P - mean(P)\n"
        "Q0 = Q - mean(Q)\n"
        "C = P0.T @ Q0\n"
        "V, S, W = svd(C)\n"
        "if det(V @ W) < 0: V[:, -1] *= -1\n"
        "R = V @ W\n"
        "P_aligned = P0 @ R + mean(Q)",
        10.0,
    )
    add_callout(slide, 7.95, 4.75, 4.55, 0.82, "This is a rigid transform, so it changes coordinate frame but not reactant fragment bond lengths.", COLORS["teal"], 11.3)
    finish(slide)

    # Slide 16
    slide = blank_slide(prs)
    add_title(slide, "Part 1: Data Contract and Target Semantics", "Technical Rationale")
    add_bullets(
        slide,
        [
            "The code uses complete reactant/product/TS triplets because supervised TS learning needs endpoint inputs and a true TS target for the same reaction ID.",
            "Atom ordering is treated as already matched across R, P, and TS. This keeps tensor construction simple, but it makes atom mapping an input contract.",
            "The current Ea target is E_TS minus the higher-energy endpoint, not always the forward barrier from the reactant. This avoids requiring an explicit reaction direction label, but the semantics must be stated clearly.",
            "Negative barriers are skipped when skip_negative_ea is enabled because they usually signal extraction, labeling, or endpoint-definition problems for this training objective.",
            "Only C/H/N/O are present in the evaluated checkpoint, so the trained atom vocabulary is intentionally narrow.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.6,
        gap=2.4,
    )
    rows = [
        ["Triplet grouping", "Require R, P, and TS logs", "Gives one supervised sample per reaction"],
        ["Atom order", "Reuse TS order for tensors", "Avoids solving atom mapping inside training"],
        ["Ea target", "E_TS - max(E_R,E_P)", "Direction-free high-side barrier used by this run"],
        ["Filtering", "max_atoms and negative Ea checks", "Keeps padded tensors bounded and labels sane"],
        ["Vocabulary", "/".join(metrics["atom_types"]), "Checkpoint cannot handle unseen atom types"],
    ]
    add_table(slide, 7.85, 1.45, 4.7, 3.2, ["Choice", "Implementation", "Why"], rows, 8.2)
    add_callout(slide, 7.85, 5.1, 4.7, 0.82, "This is the slide to use when someone asks what the labels mean and what assumptions the dataset makes.", COLORS["orange"], 11.0)
    finish(slide)

    # Slide 17
    slide = blank_slide(prs)
    add_title(slide, "Part 2: Why Distances, Masks, and Fragments", "Technical Rationale")
    add_bullets(
        slide,
        [
            "The supervised geometry target is D_TS, a pairwise distance matrix, because distances do not change under global translation or rotation.",
            "The 30-atom padded tensor shape keeps batching simple; atom masks prevent padded rows and columns from contributing to attention, loss, or metrics.",
            "The geometry mask is derived from TS fragments so disconnected cross-fragment distances do not dominate the loss with arbitrary relative placement.",
            "Fragment-aware processing matters for association and dissociation reactions: one global rigid alignment can create a false midpoint for separated species.",
            "D_I is a midpoint prior between endpoints, giving the model a chemically plausible starting distance scale instead of asking it to predict from scratch.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.6,
        gap=2.4,
    )
    add_code_box(
        slide,
        7.9,
        1.45,
        4.55,
        1.35,
        "D_R  = pairwise(c_R)\n"
        "D_P  = pairwise(c_P)\n"
        "D_I  = 0.5 * (D_R + D_P)\n"
        "loss = loss * geom_mask * valid_atom_mask",
        9.4,
    )
    rows = [
        ["Distance target", "Removes frame dependence"],
        ["Atom mask", "Ignores zero-padding"],
        ["Fragment mask", "Ignores arbitrary disconnected placement"],
        ["Midpoint prior", "Stabilizes early geometry prediction"],
    ]
    add_table(slide, 7.9, 3.35, 4.55, 1.75, ["Mechanism", "Reason"], rows, 9.1)
    add_callout(slide, 7.9, 5.55, 4.55, 0.66, "The model learns reaction geometry, not coordinate-frame bookkeeping.", COLORS["teal"], 11.5)
    finish(slide)

    # Slide 18
    slide = blank_slide(prs)
    add_title(slide, "Part 3: Why the R -> I -> P Encoder", "Technical Rationale")
    add_bullets(
        slide,
        [
            "For each atom, the model reads three distance-based states: reactant, midpoint, and product.",
            "Gaussian radial basis features turn raw distances into smooth local channels, so small distance changes produce smooth feature changes.",
            "The bidirectional GRU treats R/I/P as a short reaction-coordinate sequence and extracts the midpoint context after seeing both endpoints.",
            "Transformer layers then let all atoms exchange information globally, which is needed because a bond change in one region can shift the TS geometry elsewhere.",
            "Pre-norm residual Transformer layers make deep attention blocks easier to optimize under long training.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.6,
        gap=2.4,
    )
    add_flow_box(slide, 8.0, 1.5, 4.25, 0.72, "Gaussian RBF", "32 channels over 0.4-6.0 A", COLORS["blue"])
    add_flow_box(slide, 8.0, 2.42, 4.25, 0.72, "BiGRU", "Reads R, I, and P as a 3-step path", COLORS["teal"])
    add_flow_box(slide, 8.0, 3.34, 4.25, 0.72, "Transformer", "Global atom-to-atom context with padding mask", COLORS["purple"])
    add_code_box(
        slide,
        8.0,
        4.45,
        4.25,
        1.05,
        "emb = [RBF(D_R), RBF(D_I), RBF(D_P)]\n"
        "h_mid = BiGRU(emb)[:, midpoint]\n"
        "node = Transformer(h_mid)",
        9.0,
    )
    add_callout(slide, 8.0, 5.85, 4.25, 0.48, "This is why the architecture sees both endpoints before predicting the TS.", COLORS["blue"], 10.6)
    finish(slide)

    # Slide 19
    slide = blank_slide(prs)
    add_title(slide, "Part 4: Why Alpha Interpolation Plus Delta", "Technical Rationale")
    add_bullets(
        slide,
        [
            "The geometry head does not directly emit arbitrary distances. It starts from an endpoint interpolation and learns a bounded correction.",
            "alpha is passed through sigmoid, so each pair starts between reactant-like and product-like distances.",
            "The final layer is initialized to zero, making alpha=0.5 and delta=0 at initialization. The first prediction is therefore the endpoint midpoint.",
            "delta is clamped to +/-3.0 A so a single bad pairwise output cannot create a huge unphysical distance.",
            "A coarse auxiliary loss supervises this pre-EGNN distance matrix directly, which gives MDS and EGNN a better seed.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.6,
        gap=2.4,
    )
    add_code_box(
        slide,
        7.9,
        1.48,
        4.6,
        1.75,
        "alpha = sigmoid(alpha_logit)\n"
        "delta = clamp(delta, -3.0, +3.0)\n"
        "D_base = alpha*D_R + (1-alpha)*D_P\n"
        "D_TS   = symmetrize(clamp(D_base + delta, 0))",
        9.1,
    )
    rows = [
        ["Midpoint init", "Stable first epoch behavior"],
        ["Sigmoid alpha", "Pairwise TS position is endpoint-bounded"],
        ["Clamped delta", "Allows chemistry-specific correction without explosions"],
        ["Symmetry/diagonal", "Keeps output a valid distance matrix shape"],
    ]
    add_table(slide, 7.9, 3.72, 4.6, 1.78, ["Design", "Why"], rows, 8.8)
    add_callout(slide, 7.9, 5.9, 4.6, 0.48, "This turns geometry prediction into correction of a physical prior.", COLORS["orange"], 10.8)
    finish(slide)

    # Slide 20
    slide = blank_slide(prs)
    add_title(slide, "Part 5: Why MDS Then EGNN", "Technical Rationale")
    add_bullets(
        slide,
        [
            "MDS converts the predicted distance matrix into a concrete 3D seed, which is necessary before a coordinate refiner can operate.",
            "The MDS seed is detached in the forward pass. This avoids relying on the unstable backward path through eigendecomposition.",
            "The code runs eigendecomposition on CPU with float64 and diagonal jitter because padded molecules create repeated zero eigenvalues.",
            "EGNN refinement then updates coordinates with E(n)-equivariant message passing, so rotations and translations are handled consistently.",
            "The final distance loss is computed after EGNN refinement, meaning the supervised geometry is the actual refined TS structure.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.6,
        gap=2.4,
    )
    add_flow_box(slide, 8.0, 1.48, 4.25, 0.72, "Coarse D_TS", "Distance head output", COLORS["blue"])
    add_flow_box(slide, 8.0, 2.34, 4.25, 0.72, "Masked MDS", "Double-center D^2; keep top 3 components", COLORS["teal"])
    add_flow_box(slide, 8.0, 3.2, 4.25, 0.72, "EGNN", "4 equivariant layers, clipped coordinate updates", COLORS["purple"])
    add_flow_box(slide, 8.0, 4.06, 4.25, 0.72, "Refined D_TS", "Recomputed from coordinates for loss and Ea", COLORS["orange"])
    add_callout(slide, 8.0, 5.2, 4.25, 0.92, "MDS provides a coordinate seed; EGNN makes chemically local adjustments without breaking rigid-motion consistency.", COLORS["teal"], 11.0)
    finish(slide)

    # Slide 21
    slide = blank_slide(prs)
    add_title(slide, "Part 6: Why the Energy Head Is Warmed Up", "Technical Rationale")
    add_bullets(
        slide,
        [
            "The learned Ea head depends on EGNN-refined TS features, so early geometry mistakes would otherwise dominate the shared representation.",
            "The training loop monitors Ea every epoch but adds Ea loss only after the warmup period.",
            "Ea is z-scored for training because raw kcal/mol targets are much larger than geometry losses and would otherwise dominate optimization.",
            "The head also receives signed reaction energy and a 20D descriptor vector so it can learn BEP-like trends and composition/displacement context.",
            "The Marcus/Hammond/OLS calculation remains a baseline, giving a sanity check against a simpler chemistry-inspired model.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.6,
        gap=2.4,
    )
    rows = [
        ["Geometry warmup", f"{metric_cfg(metrics, 'ea_warmup_epochs', metric_cfg(metrics, 'energy_ramp_epochs', 200))} epochs", "Avoids learning energy from poor early TS guesses"],
        ["Ea target", "z-scored kcal/mol", "Keeps scale compatible with geometry losses"],
        ["Inputs", "pooled EGNN + de_rxn + 20D features", "Combines local TS environment with global reaction context"],
        ["Baseline", "Marcus/Hammond/OLS", "Tests whether learned head adds value"],
    ]
    add_table(slide, 7.85, 1.5, 4.7, 2.4, ["Part", "Implementation", "Why"], rows, 8.0)
    add_code_box(
        slide,
        7.85,
        4.35,
        4.7,
        1.2,
        "if epoch > ea_warmup:\n"
        "    loss += w_ea * SmoothL1(Ea_norm)\n\n"
        "Ea = Ea_norm * ea_std + ea_mean",
        9.0,
    )
    add_callout(slide, 7.85, 5.88, 4.7, 0.48, "The head is deliberately downstream of geometry quality.", COLORS["red"], 10.8)
    finish(slide)

    # Slide 22
    slide = blank_slide(prs)
    add_title(slide, "Part 7: Why the Physics Loss Terms Exist", "Technical Rationale")
    add_bullets(
        slide,
        [
            "Huber geometry loss is less sensitive than MSE to a few difficult atom pairs while still rewarding small-distance accuracy.",
            "The coarse loss trains the pairwise head even though the final loss is measured after EGNN refinement.",
            "Spectator loss encodes the assumption that pairs with little endpoint change should stay near the endpoint midpoint.",
            "Active-pair bounds allow bond-forming and bond-breaking distances to move, but penalize values far outside the endpoint range plus margin.",
            "Gradient clipping keeps rare bad batches from producing a large optimizer step during long training.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.6,
        gap=2.4,
    )
    add_code_box(
        slide,
        7.9,
        1.45,
        4.55,
        2.15,
        "L = Huber(D_refined, D_true)\n"
        "  + w_coarse * Huber(D_coarse, D_true)\n"
        "  + 0.2 * (L_spectator + 0.5*L_bounds)\n\n"
        "clip_grad_norm_(model.parameters(), 1.0)",
        8.9,
    )
    rows = [
        ["spectator_threshold", f"{metric_cfg(metrics, 'spectator_threshold', 0.15)} A"],
        ["bounds margin", "0.2 A"],
        ["geom_coarse_weight", metric_cfg(metrics, "geom_coarse_weight", 0.5)],
        ["grad_clip", metric_cfg(metrics, "grad_clip", 1.0)],
    ]
    add_table(slide, 7.9, 4.05, 4.55, 1.45, ["Setting", "Value"], rows, 9.0)
    add_callout(slide, 7.9, 5.85, 4.55, 0.48, "The loss mixes supervised data fit with simple chemistry constraints.", COLORS["purple"], 10.8)
    finish(slide)

    # Slide 23
    slide = blank_slide(prs)
    add_title(slide, "Part 8: Artifact Contract and Reproducibility", "Technical Rationale")
    add_bullets(
        slide,
        [
            "The checkpoint is more than weights: it stores atom vocabulary and normalization statistics that inference must reuse exactly.",
            "Training/validation splits use split_seed, and normalization statistics are computed only from train indices to avoid validation leakage.",
            "Coordinate augmentation improves robustness, but the current dataset noise path uses random sampling and should be seeded per worker for exact reproducibility.",
            "The current source expects physics_ea_coeffs for prediction-time physics baseline; older checkpoints may not contain that metadata and should be regenerated or guarded.",
            "Post-processing should be centralized because training evaluation and prediction currently apply slightly different constraint sequences.",
        ],
        0.72,
        1.42,
        6.9,
        5.05,
        12.4,
        gap=2.4,
    )
    rows = [
        ["atom_vocab", "Maps atom symbols to model IDs"],
        ["aphys/efeat stats", "Recreates training normalization at inference"],
        ["ea_mean/ea_std", "Denormalizes learned Ea output"],
        ["config_snapshot", "Records key architecture and training settings"],
        ["physics_ea_coeffs", "Needed by current predict baseline path"],
    ]
    add_table(slide, 7.85, 1.45, 4.7, 2.7, ["Artifact", "Why it matters"], rows, 8.4)
    add_callout(slide, 7.85, 4.65, 4.7, 0.92, "Practical rule: a prediction run must use the same feature construction, normalization, and post-processing contract as training.", COLORS["red"], 11.0)
    add_callout(slide, 7.85, 5.85, 4.7, 0.48, "This is the main engineering risk to fix before benchmark claims.", COLORS["orange"], 10.8)
    finish(slide)

    # Slide 24
    slide = blank_slide(prs)
    add_title(slide, "Known Limits and Next Technical Work", "Assessment")
    rows = [
        ["Generalization", "Train/validation gap remains large", "Stratify by reaction family and molecular size"],
        ["Target definition", "Ea is TS minus higher-energy endpoint", "Clarify forward/reverse barrier semantics"],
        ["Inference assumption", "Requires matched atom ordering", "Add atom-mapping validation or mapping step"],
        ["Chemistry scope", f"Observed atom types: {'/'.join(metrics['atom_types'])}", "Expand data and validate larger vocabularies"],
        ["Post-processing drift", "Training/eval/inference constraints differ", "Centralize one shared post-processing function"],
        ["Calibration", "Negative Ea predictions can occur", "Add calibrated non-negative handling if required"],
    ]
    add_table(slide, 0.75, 1.42, 11.85, 4.4, ["Area", "Current behavior", "Technique-focused next step"], rows, 9.6)
    add_callout(slide, 1.05, 6.15, 11.15, 0.62, "The next highest-value work is not another summary figure; it is correctness alignment, stratified validation, and uncertainty/calibration around the learned Ea head.", COLORS["red"], 12)
    finish(slide)


def main():
    metrics = collect_metrics()
    charts = generate_charts(metrics)
    prs = init_prs()
    add_algorithm_focused_deck_slides(prs, metrics, charts)

    # Ensure all text frames shrink if an installed font differs from the authoring environment.
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
